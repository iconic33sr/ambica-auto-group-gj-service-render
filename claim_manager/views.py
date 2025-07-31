from django.shortcuts import render, redirect
from django.contrib.auth.decorators import login_required
from django.contrib.auth import logout
from django.contrib import messages
from core.models import Customer_Information_Report, Service_Advisor_Report, Claim_Status, Accounting_Year
from core.forms import MyProfileForm, Manual_Password_Change_Form
from .forms import Claim_Status_Form
from datetime import date, datetime
from django.utils import timezone
from django.urls import reverse
from django.http import JsonResponse
from django.core.paginator import Paginator
from core.decorators import login_active_user_required


# For Powerpoint Presentation
import uuid
import os
import io
from django.db import models
from django.core.files.base import ContentFile
from django.shortcuts import get_object_or_404
from django.http import HttpResponse, Http404
from django.conf import settings # Make sure settings.py has MEDIA_ROOT configured

# Import python-pptx libraries
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from PIL import Image # Required for image dimension/DPI reading

from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR
import textwrap

####################################################################################################################

@login_active_user_required
def claim_manager_pending_cir_list(request):
    if request.user.user_profile.user_designation.designation == "claim_manager":
        cir_reports = Customer_Information_Report.objects.filter(workshop_manager_verification="verified", presentation_report_status="pending", claim_manager_rejection="", branch=request.user.user_profile.user_branch, selected_claim_manager=request.user.username).values('cir_uid', 'job_no', 'vehicle_no', 'supervisor_name', 'cir_date_time', 'service_advisor_report__advisor_name', 'service_advisor_report__sar_date_time', 'workshop_manager_name', 'workshop_manager_verification_date_time', 'claim_manager_preview').order_by('cir_date_time')
        paginator = Paginator(cir_reports, 30)
        page_no = request.GET.get('page')
        page_obj = paginator.get_page(page_no)
        return render(request, 'claim_manager/cir_list.html', {'cir_reports':page_obj})
    
    else:
        logout(request)
        messages.error(request, "Unauthorised Access !!")
        return redirect('user_login')
    
    
####################################################################################################################

@login_active_user_required
def claim_manager_completed_cir_list(request):
    if request.user.user_profile.user_designation.designation == "claim_manager":
        cir_reports = Customer_Information_Report.objects.filter(presentation_report_status="created", claim_manager_rejection="", branch=request.user.user_profile.user_branch, selected_claim_manager=request.user.username).values('cir_uid', 'job_no', 'vehicle_no', 'supervisor_name', 'cir_date_time', 'service_advisor_report__advisor_name', 'service_advisor_report__sar_date_time', 'workshop_manager_name', 'workshop_manager_verification_date_time', 'claim_manager_preview').order_by('cir_date_time')
        paginator = Paginator(cir_reports, 30)
        page_no = request.GET.get('page')
        page_obj = paginator.get_page(page_no)
        return render(request, 'claim_manager/completed_cir_list.html', {'cir_reports':page_obj})
    
    else:
        logout(request)
        messages.error(request, "Unauthorised Access !!")
        return redirect('user_login')
    

####################################################################################################################

@login_active_user_required
def claim_manager_rejected_cir_list(request):
    if request.user.user_profile.user_designation.designation == "claim_manager":
        cir_reports = Customer_Information_Report.objects.filter(claim_manager_rejection="rejected", branch=request.user.user_profile.user_branch, selected_claim_manager=request.user.username).values('cir_uid', 'job_no', 'vehicle_no', 'supervisor_name', 'cir_date_time', 'service_advisor_report__advisor_name', 'service_advisor_report__sar_date_time', 'workshop_manager_name', 'workshop_manager_verification_date_time', 'claim_manager_preview').order_by('cir_date_time')
        paginator = Paginator(cir_reports, 30)
        page_no = request.GET.get('page')
        page_obj = paginator.get_page(page_no)
        return render(request, 'claim_manager/rejected_cir_list.html', {'cir_reports':page_obj})
    
    else:
        logout(request)
        messages.error(request, "Unauthorised Access !!")
        return redirect('user_login')
    

####################################################################################################################

@login_active_user_required
def claim_manager_preview_cir(request, cir_uid):
    if request.user.user_profile.user_designation.designation == "claim_manager":
        if request.method == "POST":
            data = request.POST

            if data:

                report = Customer_Information_Report.objects.filter(cir_uid = data.get('cir_uid')).first()
                sar_report = Service_Advisor_Report.objects.filter(cir = report).first()

                if report.selected_claim_manager == request.user.username:

                    if report:

                        if data.get('job_no'):
                            report.job_no=data.get('job_no')
                        if data.get('vehicle_no'):
                            report.vehicle_no=data.get('vehicle_no').lower()
                        if data.get('kilometer'):
                            report.kilometer=data.get('kilometer')

                        if data.get('chassis_no'):
                            report.chassis_no=data.get('chassis_no').lower()
                        else:
                            if data.get('task') != "generate":
                                report.chassis_no=''

                        if data.get('model'):
                            report.model=data.get('model').lower()
                        else:
                            if data.get('task') != "generate":
                                report.model=''

                        if data.get('hours') in [None, '']:
                            report.hours = None
                        else:
                            report.hours = data.get('hours')

                        if data.get('sale_date'):
                            report.sale_date=data.get('sale_date')
                        else:
                            if data.get('task') != "generate":
                                report.sale_date=None

                        if data.get('complaint_1'):
                            report.complaint_1=data.get('complaint_1').lower()
                        else:
                            report.complaint_1=''
                        if data.get('complaint_2'):
                            report.complaint_2=data.get('complaint_2').lower()
                        else:
                            report.complaint_2=''
                        if data.get('complaint_3'):
                            report.complaint_3=data.get('complaint_3').lower()
                        else:
                            report.complaint_3=''
                        if data.get('complaint_4'):                
                            report.complaint_4=data.get('complaint_4').lower()
                        else:
                            report.complaint_4=''
                        if data.get('complaint_5'):                
                            report.complaint_5=data.get('complaint_5').lower()
                        else:
                            report.complaint_5=''
                        if data.get('complaint_6'):                
                            report.complaint_6=data.get('complaint_6').lower()
                        else:
                            report.complaint_6=''
                        if data.get('complaint_7'):                
                            report.complaint_7=data.get('complaint_7').lower()
                        else:
                            report.complaint_7=''

                        if data.get('workshop_manager_remark'):
                            report.workshop_manager_remark=data.get('workshop_manager_remark').lower()
                        else:
                            report.workshop_manager_remark=''

                        if data.get('claim_manager_investigation'):
                            report.claim_manager_investigation=data.get('claim_manager_investigation').lower()
                        else:
                            report.claim_manager_investigation=''
                        if data.get('claim_manager_action_taken'):
                            report.claim_manager_action_taken=data.get('claim_manager_action_taken').lower()
                        else:
                            report.claim_manager_action_taken=''

                        if data.get('task') == "reject":
                            report.claim_manager_rejection='rejected'
                            report.claim_manager_rejection_reason=data.get('claim_manager_rejection_reason')

                        report.claim_manager_last_save_date_time=datetime.now()
                        report.claim_manager_name=request.user.first_name+ " " + request.user.last_name
                        report.claim_manager_id=request.user.username
                        
                        if data.get('job_no') and data.get('vehicle_no') and data.get('kilometer'):
                            report.save()
                        else:
                            messages.error(request, "Required fields can't be empty !!")
                            return redirect('claim_manager_preview_cir', cir_uid=report.cir_uid)


                        if data.get('first_service_remark'):
                            sar_report.first_service_remark=data.get('first_service_remark').lower()
                        else:
                            sar_report.first_service_remark=''
                        if data.get('second_service_remark'):
                            sar_report.second_service_remark=data.get('second_service_remark').lower()
                        else:
                            sar_report.second_service_remark=''
                        if data.get('third_service_remark'):
                            sar_report.third_service_remark=data.get('third_service_remark').lower()
                        else:
                            sar_report.third_service_remark=''
                        if data.get('fourth_service_remark'):
                            sar_report.fourth_service_remark=data.get('fourth_service_remark').lower()
                        else:
                            sar_report.fourth_service_remark=''
                        if data.get('fifth_service_remark'):
                            sar_report.fifth_service_remark=data.get('fifth_service_remark').lower()
                        else:
                            sar_report.fifth_service_remark=''
                        if data.get('sixth_service_remark'):
                            sar_report.sixth_service_remark=data.get('sixth_service_remark').lower()
                        else:
                            sar_report.sixth_service_remark=''

                        if data.get('fault_1'):
                            sar_report.faulty1_description=data.get('fault_1').lower()
                        else:
                            sar_report.faulty1_description=''
                        if data.get('fault_2'):                
                            sar_report.faulty2_description=data.get('fault_2').lower()
                        else:
                            sar_report.faulty2_description=''
                        if data.get('fault_3'):
                            sar_report.faulty3_description=data.get('fault_3').lower()
                        else:
                            sar_report.faulty3_description=''
                        if data.get('fault_4'):
                            sar_report.faulty4_description=data.get('fault_4').lower()
                        else:
                            sar_report.faulty4_description=''
                        if data.get('fault_5'):
                            sar_report.faulty5_description=data.get('fault_5').lower()
                        else:
                            sar_report.faulty5_description=''
                        if data.get('fault_6'):
                            sar_report.faulty6_description=data.get('fault_6').lower()
                        else:
                            sar_report.faulty6_description=''
                        if data.get('fault_7'):
                            sar_report.faulty7_description=data.get('fault_7').lower()
                        else:
                            sar_report.faulty7_description=''
                        if data.get('fault_8'):
                            sar_report.faulty8_description=data.get('fault_8').lower()
                        else:
                            sar_report.faulty8_description=''
                        if data.get('fault_9'):
                            sar_report.faulty9_description=data.get('fault_9').lower()
                        else:
                            sar_report.faulty9_description=''
                        if data.get('fault_10'):
                            sar_report.faulty10_description=data.get('fault_10').lower()
                        else:
                            sar_report.faulty10_description=''
                        if data.get('fault_11'):
                            sar_report.faulty11_description=data.get('fault_11').lower()
                        else:
                            sar_report.faulty11_description=''
                        if data.get('fault_12'):
                            sar_report.faulty12_description=data.get('fault_12').lower()
                        else:
                            sar_report.faulty12_description=''
                        if data.get('fault_13'):
                            sar_report.faulty13_description=data.get('fault_13').lower()
                        else:
                            sar_report.faulty13_description=''
                        if data.get('fault_14'):
                            sar_report.faulty14_description=data.get('fault_14').lower()
                        else:
                            sar_report.faulty14_description=''
                        if data.get('fault_15'):
                            sar_report.faulty15_description=data.get('fault_15').lower()
                        else:
                            sar_report.faulty15_description=''

                        if data.get('service_advisor_action_remark'):
                            sar_report.action_remark=data.get('service_advisor_action_remark').lower()
                        else:
                            sar_report.action_remark=''
                        if data.get('advisor_remark'):                
                            sar_report.advisor_description=data.get('advisor_remark').lower()
                        else:
                            sar_report.advisor_description=''

                        sar_report.save()
                    
                        if data.get('task') == "save":
                            messages.success(request, "Data saved successfully")
                            return redirect('claim_manager_preview_cir', cir_uid=report.cir_uid)
                        
                        elif data.get('task') == "reject":
                            messages.success(request, "CIR Rejected")
                            return redirect('claim_manager_cir_list')
                        
                        else:
                            if data.get('chassis_no') != '' and data.get('model') != '' and data.get('sale_date') != '':
                                return redirect('claim_manager_generate_ppt', cir_uid=report.cir_uid)
                            else:
                                messages.error(request, "Fill the required fields !!")
                                return redirect('claim_manager_preview_cir', cir_uid=report.cir_uid)
                        
                    else:
                        messages.error(request, "Error !!")
                        return redirect('claim_manager_preview_cir', cir_uid=data.get('cir_uid'))        
                    
                else:
                    messages.error(request, "Unauthorised Access !!")
                    return redirect('user_login')                
                
            else:
                messages.error(request, "Invalid Data !!")
                return redirect('claim_manager_cir_list')                
 

        else:

            cir_report = Customer_Information_Report.objects.get(cir_uid = cir_uid)
            if cir_report.selected_claim_manager == request.user.username:
            
                try:
                    if cir_report.claim_manager_preview == "pending":
                        cir_report.claim_manager_preview = "previewed"
                        cir_report.save(update_fields=['claim_manager_preview'])
                except:
                    pass
                return render(request, 'claim_manager/preview_cir.html', {'cir_report':cir_report})    
            
            else:
                logout(request)
                messages.error(request, "Unauthorised Access !!")
                return redirect('user_login')                

    else:
        logout(request)
        messages.error(request, "Unauthorised Access !!")
        return redirect('user_login')        


####################################################################################################################

@login_active_user_required
def claim_manager_preview_completed_cir(request, cir_uid):
    if request.user.user_profile.user_designation.designation == "claim_manager":
        
        cir_report = Customer_Information_Report.objects.filter(cir_uid = cir_uid).first()
        
        if cir_report:
            return render(request, 'claim_manager/preview_completed_cir.html', {'cir_report':cir_report})    
        else:
            messages.error(request, "Error in fetching Report !!")
            return redirect('user_login')        

    else:
        logout(request)
        messages.error(request, "Unauthorised Access !!")
        return redirect('user_login')        
    

####################################################################################################################

@login_active_user_required
def claim_manager_preview_rejected_cir(request, cir_uid):
    if request.user.user_profile.user_designation.designation == "claim_manager":
            
        if request.method == "POST":
            data = request.POST

            if data:
                report = Customer_Information_Report.objects.filter(cir_uid = data.get('cir_uid')).first()

                if report:
                    if report.selected_claim_manager == request.user.username:

                        report.claim_manager_rejection = ""
                        report.claim_manager_rejection_reason = ""
                        report.save()

                        messages.success(request, "Report Unrejected Successfully. Please view it in Pending CIR.")
                        return redirect('claim_manager_rejected_cir_list')            

                    else:
                        logout(request)
                        messages.error(request, "Unauthorised Access !!")
                        return redirect('user_login')        
                    
                else:
                    messages.error(request, "Report Not Found !!")
                    return redirect('claim_manager_rejected_cir_list')            
                
            else:
                messages.error(request, "Invalid Data !!")
                return redirect('claim_manager_rejected_cir_list')        

        else:
            cir_report = Customer_Information_Report.objects.filter(cir_uid = cir_uid).first()
            
            if cir_report:
                return render(request, 'claim_manager/preview_rejected_cir.html', {'cir_report':cir_report})    
            else:
                messages.error(request, "Error in fetching Report !!")
                return redirect('user_login')

    else:
        logout(request)
        messages.error(request, "Unauthorised Access !!")
        return redirect('user_login')        


####################################################################################################################

@login_active_user_required
def claim_manager_generate_ppt(request, cir_uid):
    if request.user.user_profile.user_designation.designation == "claim_manager":

        if request.method == "POST":

            data = request.POST
                
            cir_report = Customer_Information_Report.objects.filter(cir_uid=data.get('cir_uid')).first()
            sar_report = Service_Advisor_Report.objects.filter(cir=cir_report).first()

            if cir_report.selected_claim_manager == request.user.username:

                selected_images = []

                if data.getlist('selected_options[]'):
                    selected_images = data.getlist('selected_options[]')

                if cir_report:

                    try:

                        cir_report.presentation_report_status="created"
                        cir_report.claim_manager_last_save_date_time=datetime.now()
                        cir_report.presentation_date_time=datetime.now()
                        cir_report.presentation_date=date.today()
                        cir_report.presentation_time=timezone.now()
                        cir_report.claim_manager_name=request.user.first_name+ " " + request.user.last_name
                        cir_report.claim_manager_id=request.user.username

                        cir_report.save()

                        branch_name = cir_report.branch.branch.strip().capitalize()

                        # For Powerpoint Presentation

                        prs = Presentation()

                        # Define standard slide dimensions (e.g., 16:9 aspect ratio)
                        # These are the overall slide dimensions, not the content area.
                        prs.slide_width = Inches(13.333)
                        prs.slide_height = Inches(7.5) # Standard 4:3 aspect ratio, common for presentations

                        # Define common dimensions for consistent layout within the content area
                        SLIDE_WIDTH = prs.slide_width
                        SLIDE_HEIGHT = prs.slide_height

                        # Template dimensions for content area within the slide
                        CONTENT_MARGIN = Inches(0.5) # Margin from slide edge to content area
                        CONTENT_LEFT = CONTENT_MARGIN
                        CONTENT_TOP = CONTENT_MARGIN
                        CONTENT_WIDTH = SLIDE_WIDTH - (2 * CONTENT_MARGIN)
                        CONTENT_HEIGHT = SLIDE_HEIGHT - (2 * CONTENT_MARGIN)

                        # Common font sizes
                        TITLE_FONT_SIZE = Pt(48)
                        SUBTITLE_FONT_SIZE = Pt(32)
                        SECTION_TITLE_FONT_SIZE = Pt(36)
                        BODY_FONT_SIZE = Pt(18)
                        LABEL_FONT_SIZE = Pt(14)


                        # Helper function to add a slide with the attractive template
                        def _add_slide_with_attractive_template(prs_obj, slide_layout_index=6):
                            slide = prs_obj.slides.add_slide(prs_obj.slide_layouts[slide_layout_index])
                            slide.background.fill.solid()
                            slide.background.fill.fore_color.rgb = RGBColor(255, 255, 255)  # White background

                            # Dimensions
                            SLIDE_WIDTH = prs_obj.slide_width
                            SLIDE_HEIGHT = prs_obj.slide_height
                            header_height = Inches(0.5)
                            footer_height = Inches(0.3)

                            # --- HEADER BAR ---
                            header_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_WIDTH, header_height)
                            header_shape.fill.solid()
                            header_shape.fill.fore_color.rgb = RGBColor(7, 7, 191)  # Dark blue
                            header_shape.line.fill.background()

                            # --- TATA LOGO LEFT ---
                            tata_logo_path = os.path.join(settings.BASE_DIR, 'core', 'static', 'core', 'images', 'tata_logo_blue.png')
                            if os.path.exists(tata_logo_path):
                                slide.shapes.add_picture(tata_logo_path, Inches(0.2), Inches(0.02), height=header_height - Inches(0.05))

                            # --- AMBICA LOGO CENTER ---
                            ambica_logo_path = os.path.join(settings.BASE_DIR, 'core', 'static', 'core', 'images', 'ambica_logo_blue.png')
                            if os.path.exists(ambica_logo_path):
                                slide.shapes.add_picture(ambica_logo_path, (SLIDE_WIDTH - Inches(0.7)) / 2, Inches(0.02), height=header_height - Inches(0.05))

                            # --- TATA MOTORS BRANDING IMAGE RIGHT ---
                            tata_motors_branding_path = os.path.join(settings.BASE_DIR, 'core', 'static', 'core', 'images', 'tata_motors_branding_blue.png')
                            if os.path.exists(tata_motors_branding_path):
                                branding_img_width = Inches(1.45)
                                branding_img_height = Inches(0.15)
                                branding_img_top = (header_height - branding_img_height) / 2

                                slide.shapes.add_picture(
                                    tata_motors_branding_path,
                                    SLIDE_WIDTH - branding_img_width - Inches(0.2),  # Fully right-aligned with 0.2in padding
                                    branding_img_top,
                                    width=branding_img_width,
                                    height=branding_img_height
                                )


                            # --- FOOTER BAR ---
                            footer_shape = slide.shapes.add_shape(
                                MSO_SHAPE.RECTANGLE, 0, SLIDE_HEIGHT - footer_height, SLIDE_WIDTH, footer_height
                            )
                            footer_shape.fill.solid()
                            footer_shape.fill.fore_color.rgb = RGBColor(159, 0, 0)  # Dark red
                            footer_shape.line.fill.background()

                            # Footer text
                            footer_text = f"Shree Ambica Auto Sales & Service (2000480) – {branch_name}"
                            footer_text_box = slide.shapes.add_textbox(
                                Inches(0.2),
                                SLIDE_HEIGHT - footer_height,
                                SLIDE_WIDTH - Inches(0.4),
                                footer_height
                            )

                            footer_frame = footer_text_box.text_frame
                            footer_frame.clear()
                            footer_frame.vertical_anchor = MSO_ANCHOR.MIDDLE  # Vertically center the text
                            p = footer_frame.paragraphs[0]
                            p.text = footer_text
                            p.font.name = 'Segoe UI Semibold'
                            p.font.size = Pt(12)
                            p.font.bold = True
                            p.font.color.rgb = RGBColor(255, 255, 255)
                            p.alignment = PP_ALIGN.CENTER

                            return slide
                        

                        def _add_subheading_box(slide_obj, heading_text):
                            subheading_top = Inches(0.6)  # Just below the header bar   
                            subheading_height = Inches(0.4)
                            subheading_width = CONTENT_WIDTH
                            subheading_left = CONTENT_LEFT

                            shape = slide_obj.shapes.add_shape(
                                MSO_SHAPE.RECTANGLE,
                                subheading_left,
                                subheading_top,
                                subheading_width,
                                subheading_height
                            )
                            shape.fill.solid()
                            shape.fill.fore_color.rgb = RGBColor(198, 224, 180)  # Light green
                            shape.line.color.rgb = RGBColor(0, 100, 0)           # Dark green border
                            shape.line.width = Pt(2)

                            text_frame = shape.text_frame
                            text_frame.clear()
                            p = text_frame.paragraphs[0]
                            p.text = heading_text.upper()
                            p.font.name = 'Segoe UI Semibold'
                            p.font.size = Pt(20)
                            p.font.bold = True
                            p.font.color.rgb = RGBColor(0, 0, 0)
                            p.alignment = PP_ALIGN.CENTER



                        # Helper function to add and scale image
                        def _add_image_with_aspect_ratio_and_label(slide_obj, image_path, target_left, target_top, target_width, target_height, label_text):
                            """
                            Adds an image to the slide, scaling it to fit within the target box while preserving aspect ratio,
                            and adds a label below it. Includes robust error handling for missing/corrupt images.
                            """
                            if not image_path or not os.path.exists(image_path):
                                # Add placeholder text for missing images
                                error_text_box = slide_obj.shapes.add_textbox(target_left, target_top, target_width, target_height)
                                error_tf = error_text_box.text_frame
                                error_tf.clear()
                                p = error_tf.paragraphs[0]
                                p.text = f"{label_text} Not Available"
                                p.alignment = PP_ALIGN.CENTER
                                p.font.size = Pt(16)
                                p.font.bold = True
                                p.font.color.rgb = RGBColor(255, 0, 0)
                                return

                            try:
                                # Open image with PIL to get dimensions and DPI
                                img = Image.open(image_path)
                                img_width_px, img_height_px = img.size

                                # Get DPI, default to 72 if not present
                                dpi_x, dpi_y = img.info.get('dpi', (72, 72))
                                dpi_x = int(dpi_x) if isinstance(dpi_x, (int, float)) else 72
                                dpi_y = int(dpi_y) if isinstance(dpi_y, (int, float)) else 72

                                # Convert pixels to EMUs (1 inch = 914400 EMUs)
                                img_width_emu = Emu(int(img_width_px / dpi_x * 914400))
                                img_height_emu = Emu(int(img_height_px / dpi_y * 914400))

                                # Calculate scaling ratio to maintain aspect ratio
                                ratio_w = target_width.emu / img_width_emu.emu
                                ratio_h = target_height.emu / img_height_emu.emu
                                ratio = min(ratio_w, ratio_h)

                                # Final image size
                                pic_width = Emu(int(img_width_emu.emu * ratio))
                                pic_height = Emu(int(img_height_emu.emu * ratio))

                                # Center image within target box
                                pic_left = Emu(int(target_left.emu + (target_width.emu - pic_width.emu) / 2))
                                pic_top = Emu(int(target_top.emu + (target_height.emu - pic_height.emu) / 2))

                                # Add image to slide
                                pic_shape = slide_obj.shapes.add_picture(image_path, pic_left, pic_top, width=pic_width, height=pic_height)

                                # Optional: add a border
                                pic_shape.line.color.rgb = RGBColor(100, 100, 100)
                                pic_shape.line.width = Emu(9525)

                                # Add label below image
                                label_top = Emu(int(pic_top.emu + pic_height.emu + Inches(0.1).emu))
                                label_height = Inches(0.3)
                                label_box = slide_obj.shapes.add_textbox(target_left, label_top, target_width, label_height)
                                label_frame = label_box.text_frame
                                label_frame.clear()
                                p = label_frame.paragraphs[0]
                                p.text = label_text
                                p.font.size = LABEL_FONT_SIZE
                                p.font.bold = True
                                p.alignment = PP_ALIGN.CENTER

                            except Exception as e:
                                # Handle any image processing error
                                error_text_box = slide_obj.shapes.add_textbox(target_left, target_top, target_width, target_height)
                                error_tf = error_text_box.text_frame
                                error_tf.clear()
                                p = error_tf.paragraphs[0]
                                p.text = f"Error loading {label_text}: {e}"
                                p.alignment = PP_ALIGN.CENTER
                                p.font.size = Pt(12)
                                p.font.color.rgb = RGBColor(255, 0, 0)
                                print(f"Error processing image {image_path}: {e}")


                        def _add_vehicle_details_middle_content(prs, slide, cir_report, sar_report):
                            SLIDE_WIDTH = prs.slide_width
                            SLIDE_HEIGHT = prs.slide_height
                            CONTENT_MARGIN = Inches(0.5)
                            CONTENT_LEFT = CONTENT_MARGIN
                            CONTENT_TOP = CONTENT_MARGIN
                            CONTENT_WIDTH = SLIDE_WIDTH - (2 * CONTENT_MARGIN)
                            
                            heading_top = Inches(0.6)
                            heading_height = Inches(0.4)
                            details_top = heading_top + heading_height + Inches(0.36)
                            details_height = Inches(2.5) # before 2.25

                            # --- Details Grey Box ---
                            details_shape = slide.shapes.add_shape(
                                MSO_SHAPE.RECTANGLE,
                                CONTENT_LEFT,
                                details_top,
                                CONTENT_WIDTH,
                                details_height
                            )
                            details_shape.fill.solid()
                            details_shape.fill.fore_color.rgb = RGBColor(245, 245, 245)
                            details_shape.line.fill.background()  # No border

                            details_tf = details_shape.text_frame
                            details_tf.clear()
                            details_tf.word_wrap = True

                            # Prepare points: (label, value)
                            detail_points = [
                                ("REG NO", cir_report.vehicle_no.upper() or 'N/A'),
                                ("CHASSIS NO", cir_report.chassis_no.upper() or 'N/A'),
                                ("MODEL", cir_report.model.upper() if cir_report.model else 'N/A'),
                            ]
                            if cir_report.hours:
                                detail_points.append(("KM / HRS", f"{cir_report.kilometer or 'N/A'} / {cir_report.hours}"))
                            else:
                                detail_points.append(("KM", cir_report.kilometer or 'N/A'))
                            detail_points += [
                                ("SALE DATE", cir_report.sale_date.strftime('%d/%m/%Y') if cir_report.sale_date else 'N/A'),
                                ("JOB NO", cir_report.job_no or 'N/A'),
                            ]

                            for i, (point, value) in enumerate(detail_points):
                                para = details_tf.add_paragraph() if i > 0 else details_tf.paragraphs[0]
                                para.clear()
                                para.alignment = PP_ALIGN.LEFT
                                para.space_after = Pt(8)
                                para.font.size = Pt(15)
                                para.font.name = 'Arial'
                                # Bullet
                                run_bullet = para.add_run()
                                run_bullet.text = "➤  "
                                run_bullet.font.size = Pt(15)
                                run_bullet.font.bold = True
                                run_bullet.font.color.rgb = RGBColor(0, 0, 0)
                                run_bullet.font.name = 'Arial'
                                # Key (bold)
                                run_key = para.add_run()
                                run_key.text = f"{point}"
                                run_key.font.bold = True
                                run_key.font.size = Pt(15)
                                run_key.font.color.rgb = RGBColor(0, 0, 0)
                                run_key.font.name = 'Arial'
                                # Colon and space
                                run_colon = para.add_run()
                                run_colon.text = " : "
                                run_colon.font.bold = True
                                run_colon.font.size = Pt(15)
                                run_colon.font.color.rgb = RGBColor(0, 0, 0)
                                run_colon.font.name = 'Arial'
                                # Value (bold only for REG NO line)
                                run_val = para.add_run()
                                run_val.text = f"{value}"
                                run_val.font.bold = True if i == 0 else False
                                run_val.font.size = Pt(15)
                                run_val.font.color.rgb = RGBColor(0, 0, 0)
                                run_val.font.name = 'Arial'

                            # --- Purple Divider Line ---
                            line_top = details_top + details_height + Inches(0.10)
                            line_height = Inches(0.12)
                            line_shape = slide.shapes.add_shape(
                                MSO_SHAPE.RECTANGLE,
                                CONTENT_LEFT,
                                line_top,
                                CONTENT_WIDTH,
                                line_height
                            )
                            line_shape.fill.solid()
                            line_shape.fill.fore_color.rgb = RGBColor(98, 0, 153)
                            line_shape.space_after = Pt(8)


                            # --- Bottom Info Box (Supervisor etc.) ---
                            box_top = line_top + line_height + Inches(0.14)
                            box_height = Inches(1.8) # before 1.3
                            box_shape = slide.shapes.add_shape(
                                MSO_SHAPE.RECTANGLE,
                                CONTENT_LEFT,
                                box_top,
                                CONTENT_WIDTH,
                                box_height
                            )
                            box_shape.fill.solid()
                            box_shape.fill.fore_color.rgb = RGBColor(245, 245, 245)
                            box_shape.line.fill.background()

                            box_tf = box_shape.text_frame
                            box_tf.clear()
                            box_tf.word_wrap = True

                            supervisor = cir_report.supervisor_name.title() or 'N/A'
                            compl_date = cir_report.cir_date_time.strftime('%d/%m/%Y') if cir_report.cir_date_time else 'N/A'
                            advisor = getattr(sar_report, 'advisor_name', '').title() or 'N/A'
                            sub_date = sar_report.sar_date_time.strftime('%d/%m/%Y') if sar_report and sar_report.sar_date_time else 'N/A'

                            info_lines = [
                                ("Supervisor", supervisor or 'N/A'),
                                ("Complain Date", compl_date or 'N/A'),
                                ("Service Advisor", advisor or 'N/A'),
                                ("Sub Date", sub_date or 'N/A'),
                            ]

                            # for i, text in enumerate(info_lines):
                            for i, (point, value) in enumerate(info_lines):
                                
                                para = box_tf.add_paragraph() if i > 0 else box_tf.paragraphs[0]
                                para.clear()
                                para.alignment = PP_ALIGN.LEFT
                                para.space_after = Pt(16) if i == 1 else Pt(8)
                                para.font.size = Pt(15)
                                para.font.name = 'Arial'
                                # Bullet
                                run_bullet = para.add_run()
                                run_bullet.text = "➤  "
                                run_bullet.font.size = Pt(15)
                                run_bullet.font.bold = True
                                run_bullet.font.color.rgb = RGBColor(0, 0, 0)
                                run_bullet.font.name = 'Arial'
                                # Key (bold)
                                run_key = para.add_run()
                                run_key.text = f"{point}"
                                run_key.font.bold = True
                                run_key.font.size = Pt(15)
                                run_key.font.color.rgb = RGBColor(0, 0, 0)
                                run_key.font.name = 'Arial'
                                # Colon and space
                                run_colon = para.add_run()
                                run_colon.text = " : "
                                run_colon.font.bold = True
                                run_colon.font.size = Pt(15)
                                run_colon.font.color.rgb = RGBColor(0, 0, 0)
                                run_colon.font.name = 'Arial'
                                # Value 
                                run_val = para.add_run()
                                run_val.text = f"{value}"
                                run_val.font.bold = False
                                run_val.font.size = Pt(15)
                                run_val.font.color.rgb = RGBColor(0, 0, 0)
                                run_val.font.name = 'Arial'


                        def add_label_and_image_optional(
                            slide, left, top, width, height, label=None, image_path=None,
                            font_size=14, label_height=0.6
                        ):
                            """
                            Add a label (wrapped, multi-line) over the image.
                            If label or image_path is None/blank, handle gracefully.
                            """
                            label = (label or '').strip()
                            image_present = image_path and os.path.exists(image_path)
                            label_present = bool(label)

                            # -- Both present --
                            if image_present and label_present:
                                # Add label box at top
                                label_box = slide.shapes.add_textbox(left, top, width, Inches(label_height))
                                tf = label_box.text_frame
                                tf.clear()
                                tf.word_wrap = True
                                tf.vertical_anchor = MSO_ANCHOR.TOP
                                wrapped = "\n".join(textwrap.wrap(label, 35))  # Adjust char-per-line as needed
                                p = tf.paragraphs[0]
                                p.text = wrapped
                                p.font.name = 'Segoe UI Semibold'
                                p.font.size = Pt(font_size)
                                p.font.bold = True
                                p.font.color.rgb = RGBColor(255, 255, 255)
                                p.alignment = PP_ALIGN.CENTER
                                # Semi-transparent bg
                                fill = label_box.fill
                                fill.solid()
                                fill.fore_color.rgb = RGBColor(0, 0, 0)
                                fill.transparency = 0.35
                                # Reduce image height to fit label
                                img_top = top + Inches(label_height)
                                img_height = height - Inches(label_height)
                                if img_height < Inches(1): img_height = Inches(1)
                                slide.shapes.add_picture(image_path, left, img_top, width=width, height=img_height)
                            # -- Only image --
                            elif image_present and not label_present:
                                slide.shapes.add_picture(image_path, left, top, width=width, height=height)
                            # -- Only label --
                            elif label_present and not image_present:
                                label_box = slide.shapes.add_textbox(left, top, width, height)
                                tf = label_box.text_frame
                                tf.clear()
                                tf.word_wrap = True
                                tf.vertical_anchor = MSO_ANCHOR.MIDDLE
                                wrapped = "\n".join(textwrap.wrap(label, 35))
                                p = tf.paragraphs[0]
                                p.text = wrapped
                                p.font.name = 'Segoe UI Semibold'
                                p.font.size = Pt(font_size)
                                p.font.bold = True
                                p.font.color.rgb = RGBColor(0, 0, 0)
                                p.alignment = PP_ALIGN.CENTER
                                fill = label_box.fill
                                fill.solid()
                                fill.fore_color.rgb = RGBColor(235, 235, 200)
                                fill.transparency = 0.0
                            # -- Neither present --
                            else:
                                # Optional: show a placeholder or skip
                                placeholder = slide.shapes.add_textbox(left, top, width, height)
                                tf = placeholder.text_frame
                                tf.text = "No image/label"
                                tf.paragraphs[0].font.size = Pt(13)
                                tf.paragraphs[0].font.color.rgb = RGBColor(200, 0, 0)
                                tf.paragraphs[0].alignment = PP_ALIGN.CENTER


                        def _add_image_slide(prs_obj, title_text, image_data_list): ############################################################
                            def add_new_slide(prs_obj, title_text):
                                slide = _add_slide_with_attractive_template(prs_obj, 6)  # Blank layout
                                _add_subheading_box(slide, title_text)  # Subheading box below header

                                return slide

                            # Layout constants
                            image_width = Inches(4)
                            if title_text in ["Vehicle Images", "Fault Parts"]:
                                image_height = Inches(5.8)
                            else:    
                                image_height = Inches(5.8)

                            spacing_x = Inches(1.0)

                            if title_text in ["Vehicle Images", "Fault Parts"]:
                                spacing_top = Inches(0.7)   # Below title
                            else:
                                spacing_top = Inches(0.7)   # Below title

                            label_height = Inches(0.4)
                            images_per_row = 2

                            idx = 0
                            slide = add_new_slide(prs_obj, title_text)
                            while idx < len(image_data_list):
                                # How many images in this row? (1 or 2)
                                row_img_count = min(images_per_row, len(image_data_list) - idx)
                                # Total width taken by images + spacings
                                total_block_width = row_img_count * image_width + (row_img_count - 1) * spacing_x
                                # Left margin to center this block
                                block_left = CONTENT_LEFT + (CONTENT_WIDTH - total_block_width) / 2

                                for col in range(row_img_count):
                                    img_data = image_data_list[idx]
                                    left = block_left + col * (image_width + spacing_x)
                                    top = CONTENT_TOP + spacing_top
                                    label = img_data.get('label')
                                    image_path = img_data.get('path')
                                    add_label_and_image_optional(
                                        slide, left, top, image_width, image_height,
                                        label=label, image_path=image_path, font_size=14, label_height=0.6
                                    )
                                    idx += 1


                                # If more images remain, start a new slide
                                if idx < len(image_data_list):
                                    slide = add_new_slide(prs_obj, title_text)




                        def _add_text_slide(prs_obj, title_text, content_text):
                            """
                            Adds a slide with a title and a main text body.
                            """
                            if not content_text or not content_text.strip():
                                return  # Skip empty content

                            # Add a blank slide
                            slide = _add_slide_with_attractive_template(prs_obj, 6)
                            _add_subheading_box(slide, title_text)

                            # --- Main Content Text Box ---
                            text_left = CONTENT_LEFT + Inches(0.2)
                            text_top = CONTENT_TOP + Inches(0.2) + Inches(0.8) + Inches(0.2)  # Top + title height + gap
                            text_width = CONTENT_WIDTH - Inches(0.4)
                            text_height = SLIDE_HEIGHT - text_top - CONTENT_MARGIN

                            content_box = slide.shapes.add_textbox(text_left, text_top, text_width, text_height)
                            content_frame = content_box.text_frame
                            content_frame.word_wrap = True
                            content_frame.clear()

                            # Add each non-empty line as a new paragraph
                            for idx, line in enumerate(content_text.split('\n')):
                                clean_line = line.strip()
                                if not clean_line:
                                    continue
                                para = content_frame.add_paragraph() if idx > 0 else content_frame.paragraphs[0]
                                para.text = clean_line
                                para.font.size = BODY_FONT_SIZE
                                para.font.bold = False
                                para.alignment = PP_ALIGN.LEFT


                        # --- Start Presentation Generation ---

                        # Slide 1: Title Slide
                        slide = _add_slide_with_attractive_template(prs, 6)  # Use blank slide layout for full control

                        # --- Main Title ---
                        box_left = CONTENT_LEFT  # usually Inches(0.5)
                        box_top = CONTENT_TOP + Inches(1.1)   # adjust top as needed to center vertically
                        box_width = CONTENT_WIDTH
                        box_height = Inches(1.3)  # adjust as needed

                        # Add rectangle for the heading box
                        heading_box_shape = slide.shapes.add_shape(
                            MSO_SHAPE.RECTANGLE,
                            box_left,
                            box_top,
                            box_width,
                            box_height
                        )
                        heading_box_shape.fill.solid()
                        heading_box_shape.fill.fore_color.rgb = RGBColor(252, 213, 181)  # white, or light color
                        heading_box_shape.line.color.rgb = RGBColor(0, 100, 0)           # dark green or your border color
                        heading_box_shape.line.width = Pt(3)                             # adjust border thickness

                        # Now add your title text INSIDE this box (so the text overlaps the shape)
                        title_box = slide.shapes.add_textbox(
                            box_left,
                            box_top,
                            box_width,
                            box_height
                        )
                        title_tf = title_box.text_frame
                        title_tf.clear()
                        p = title_tf.paragraphs[0]
                        p.text = "COMPLAIN INFORMATION REPORT"
                        p.font.size = Pt(60)
                        p.font.bold = True
                        p.font.color.rgb = RGBColor(0, 32, 128)  # deep blue
                        p.alignment = PP_ALIGN.CENTER
                        title_tf.vertical_anchor = MSO_ANCHOR.MIDDLE


                        # Adjust font size to fit within textbox
                        max_font_size = 44  # Start from large size
                        min_font_size = 18  # Do not go below this

                        # Function to check if text fits
                        def text_fits(paragraph, textbox):
                            # Rough estimate: decrease font size if text is too long for the box width
                            text_length = len(paragraph.text)
                            box_width_inches = textbox.width.inches
                            return text_length * paragraph.font.size.pt * 0.01 < box_width_inches * 72

                        font_size = max_font_size
                        while font_size >= min_font_size:
                            p.font.size = Pt(font_size)
                            if text_fits(p, title_box):
                                break
                            font_size -= 2

                        # --- Registration Number ---
                        reg_no_box = slide.shapes.add_textbox(
                            CONTENT_LEFT,
                            CONTENT_TOP + Inches(1) + Inches(1.5) + Inches(0.5),
                            CONTENT_WIDTH,
                            Inches(0.75)
                        )
                        reg_no_tf = reg_no_box.text_frame
                        reg_no_tf.clear()
                        p = reg_no_tf.paragraphs[0]
                        p.text = f"VEHICLE NO : {cir_report.vehicle_no.upper() or 'N/A'}"
                        p.font.size = SUBTITLE_FONT_SIZE
                        p.font.bold = False
                        p.alignment = PP_ALIGN.CENTER
                        p.font.color.rgb = RGBColor(0, 0, 128)


                        # --- Slide 2: VEHICLE DETAILS ---
                        vehicle_details_slide = _add_slide_with_attractive_template(prs, 6)
                        _add_subheading_box(vehicle_details_slide, "VEHICLE DETAILS")
                        _add_vehicle_details_middle_content(prs, vehicle_details_slide, cir_report, sar_report)


                        # --- Vehicle Images ---
                        vehicle_images = []
                        if cir_report.vehicle_front_image:
                            vehicle_images.append({'path': cir_report.vehicle_front_image.path, 'label': 'VEHICLE FRONT'})
                        if cir_report.vehicle_with_number_plate:
                            vehicle_images.append({'path': cir_report.vehicle_with_number_plate.path, 'label': 'VEHICLE WITH NUMBER PLATE'})
                        if cir_report.chasis:
                            vehicle_images.append({'path': cir_report.chasis.path, 'label': 'CHASSIS'})
                        if cir_report.odometer:
                            vehicle_images.append({'path': cir_report.odometer.path, 'label': 'ODOMETER'})

                        if vehicle_images:
                            _add_image_slide(prs, "Vehicle Images", vehicle_images)

                        # --- Complaint Images ---
                        complaint_images = []
                        if cir_report.complaint_1_image or cir_report.complaint_1:
                            complaint_images.append({'path': cir_report.complaint_1_image.path if cir_report.complaint_1_image else '', 'label': cir_report.complaint_1.upper() or ''})
                        if cir_report.complaint_2_image  or cir_report.complaint_2:
                            complaint_images.append({'path': cir_report.complaint_2_image.path if cir_report.complaint_2_image else '', 'label': cir_report.complaint_2.upper() or ''})
                        if cir_report.complaint_3_image  or cir_report.complaint_3:
                            complaint_images.append({'path': cir_report.complaint_3_image.path if cir_report.complaint_3_image else '', 'label': cir_report.complaint_3.upper() or ''})
                        if cir_report.complaint_4_image  or cir_report.complaint_4:
                            complaint_images.append({'path': cir_report.complaint_4_image.path if cir_report.complaint_4_image else '', 'label': cir_report.complaint_4.upper() or ''})
                        if cir_report.complaint_5_image  or cir_report.complaint_5:
                            complaint_images.append({'path': cir_report.complaint_5_image.path if cir_report.complaint_5_image else '', 'label': cir_report.complaint_5.upper() or ''})
                        if cir_report.complaint_6_image  or cir_report.complaint_6:
                            complaint_images.append({'path': cir_report.complaint_6_image.path if cir_report.complaint_6_image else '', 'label': cir_report.complaint_6.upper() or ''})
                        if cir_report.complaint_7_image  or cir_report.complaint_7:
                            complaint_images.append({'path': cir_report.complaint_7_image.path if cir_report.complaint_7_image else '', 'label': cir_report.complaint_7.upper() or ''})

                        if complaint_images:
                            _add_image_slide(prs, "Customer Complain", complaint_images)

                        # --- Fault Images ---

                        # For both images and description
                        fault_images = []
                        if selected_images:
                            for image_name in selected_images:  # selected_images = ['faulty_image_1', 'faulty_image_2', ...]
                                # Extract index after the second underscore
                                image_name_index = image_name.split('_')[2]  # e.g. '1' from 'faulty_image_1'
                                description_field = f"faulty{image_name_index}_description"

                                # Get the image field and description field from sar_report
                                image_field = getattr(sar_report, image_name, None)
                                description_value = getattr(sar_report, description_field, '')

                                fault_images.append({
                                    'path': image_field.path if image_field else '',
                                    'label': description_value.upper() or ''
                                })


                        if fault_images:
                            _add_image_slide(prs, "Fault Parts", fault_images)

                        # --- Text-Only Slides ---
                        if cir_report.claim_manager_investigation:
                            _add_text_slide(prs, "Investigation", cir_report.claim_manager_investigation.upper())

                        if cir_report.claim_manager_action_taken:
                            _add_text_slide(prs, "Action Taken", cir_report.claim_manager_action_taken.upper())

                        # --- Save presentation to memory ---
                        ppt_buffer = io.BytesIO()
                        prs.save(ppt_buffer)
                        ppt_buffer.seek(0)

                        # --- Generate filename ---
                        base_filename = cir_report.job_no or str(cir_report.vehicle_no)
                        filename = f"{base_filename.replace(' ', '_')}_{uuid.uuid4().hex[:8]}.pptx"

                        # --- Save to FileField in model ---
                        cir_report.presentation_report.save(filename, ContentFile(ppt_buffer.getvalue()))

                        return redirect('presentation_download', cir_uid=cir_report.pk)
                    
                    except:
                        messages.error(request, "Error in generating presentation !!")    
                        return redirect('claim_manager_preview_cir', cir_uid=data.get('cir_uid'))        
                

                else:
                    messages.error(request, "Error !!")
                    return redirect('claim_manager_preview_cir', cir_uid=data.get('cir_uid'))        

            else:
                messages.error(request, "Unauthorised Access !!")
                return redirect('claim_manager_cir_list')

        else:
            cir_report = Customer_Information_Report.objects.filter(cir_uid=cir_uid).first()

            if cir_report.selected_claim_manager == request.user.username:
                sar_report = Service_Advisor_Report.objects.filter(cir=cir_report).first()
                return render(request, 'claim_manager/select_faulty_images.html', {'cir_uid':cir_report.cir_uid, 'job_no':cir_report.job_no, 'vehicle_no':cir_report.vehicle_no, 'sar_report':sar_report})
            else:
                logout(request)
                messages.error(request, "Unauthorised Access !!")
                return redirect('user_login')                    
        
    else:
        logout(request)
        messages.error(request, "Unauthorised Access !!")
        return redirect('user_login')
        

####################################################################################################################

@login_active_user_required
def presentation_download(request, cir_uid):
    if request.user.user_profile.user_designation.designation == "claim_manager":

        cir_report = get_object_or_404(Customer_Information_Report, pk=cir_uid)
        if cir_report.selected_claim_manager == request.user.username:

            file_url = cir_report.presentation_report.url
            redirect_url = reverse('claim_manager_cir_list')
            return render(request, 'claim_manager/presentation_download.html', {
                'file_url': file_url,
                'redirect_url': redirect_url
            })
        
        else:
            messages.error(request, "Unauthorised Access !!")
            return redirect('claim_manager_cir_list')    
    
    else:
        logout(request)
        messages.error(request, "Unauthorised Access !!")
        return redirect('user_login')


########################################################################################################################################

@login_active_user_required
def claim_status(request):
    if request.user.user_profile.user_designation.designation == "claim_manager":

        if request.method == 'POST':
            selected_job_no = request.POST.get('job_card_no')

            if selected_job_no:
                if selected_job_no[0:2] == "JC":

                    cform = Claim_Status_Form(request.POST)
                    if cform.is_valid():    
                        cir_report = Customer_Information_Report.objects.filter(job_no = selected_job_no).first()
                        if cir_report:
                            if cir_report.branch.branch == request.user.user_profile.user_branch.branch:
                                try:
                                    claim_status_report = cir_report.claim_status
                                except Claim_Status.DoesNotExist:
                                    claim_status_report = None

                                creating_new = claim_status_report is None

                                if creating_new:
                                    claim_status_report = Claim_Status(cir = cir_report)

                                claim_status_report.claim_no = (cform.cleaned_data.get('claim_no') or '').lstrip('0')
                                claim_status_report.claim_amount = cform.cleaned_data.get('claim_amount') or ''
                                claim_status_report.claim_date = cform.cleaned_data.get('claim_date') or None
                                claim_status_report.submission_date = cform.cleaned_data.get('submission_date') or None
                                claim_status_report.claim_status = cform.cleaned_data.get('claim_status') or 'pending'
                                claim_status_report.claim_settled_date = cform.cleaned_data.get('claim_settled_date') or None
                                if cform.cleaned_data.get('claim_status') != "rejected":
                                    claim_status_report.crm_rejection_reason = ''
                                else:
                                    claim_status_report.crm_rejection_reason = cform.cleaned_data.get('crm_rejection_reason') or ''
                                
                                claim_status_report.save()

                                messages.success(request, "Data saved successfully")    
                            else:
                                logout(request)
                                messages.error(request, "Unauthorised Access !!")
                                return redirect('user_login')    

                        else:
                            messages.error(request, "Report doesn't exist !!")    
                        
                    else:        
                        messages.error(request, "Invalid Data !!")

                else:        
                    messages.error(request, "Invalid Job No !!")

            else:
                messages.error(request, "Job No can't be empty !!")

            return redirect('claim_manager_claim_status')        

        else:
            ac_years = Accounting_Year.objects.all().order_by('-ac_year')
            cform = Claim_Status_Form()
            return render(request, 'claim_manager/claim_status.html', {'cform':cform, 'ac_years':ac_years})
        
    else:
        logout(request)
        messages.error(request, "Unauthorised Access !!")
        return redirect('user_login')    
    

########################################################

@login_active_user_required
def fetch_claim_status_data(request):
    if request.user.user_profile.user_designation.designation == "claim_manager":
        if request.method == "GET":
            full_job_no = "JC-ShrAmb-" + request.user.user_profile.user_branch.code.upper() + "-" + request.GET.get("job_no","")
            cir_report = Customer_Information_Report.objects.filter(job_no = full_job_no).first()
            if not cir_report:
                return JsonResponse({"error_msg": "Job No not found !!"})
            
            elif cir_report.presentation_report_status == "pending" and cir_report.claim_manager_rejection == "":
                return JsonResponse({"error_msg": "Report is Pending !!"})
            
            elif cir_report.claim_manager_rejection == "rejected":
                return JsonResponse({"error_msg": "Report is Rejected !!"})

            else:
            
                try:
                    claim_status = cir_report.claim_status
                except Claim_Status.DoesNotExist:
                    return JsonResponse({"vehicle_no": cir_report.vehicle_no, "full_job_no":full_job_no})
                
                data = {
                    "full_job_no": full_job_no,
                    "vehicle_no": cir_report.vehicle_no if cir_report.vehicle_no else "",
                    "claim_no": claim_status.claim_no if claim_status.claim_no else "",
                    "claim_amount": claim_status.claim_amount if claim_status.claim_amount else "",
                    "claim_date": claim_status.claim_date.strftime("%Y-%m-%d") if claim_status.claim_date else "",
                    "submission_date": claim_status.submission_date.strftime("%Y-%m-%d") if claim_status.submission_date else "",
                    "claim_settled_date": claim_status.claim_settled_date.strftime("%Y-%m-%d") if claim_status.claim_settled_date else "",
                    "claim_status": claim_status.claim_status if claim_status.claim_status else "",
                    "crm_rejection_reason": claim_status.crm_rejection_reason if claim_status.crm_rejection_reason else "",
                    "part_dispatch_image1": claim_status.part_dispatch_image1.url if claim_status.part_dispatch_image1 else "",
                    "part_dispatch_image2": claim_status.part_dispatch_image2.url if claim_status.part_dispatch_image2 else "",
                    "part_dispatch_image3": claim_status.part_dispatch_image3.url if claim_status.part_dispatch_image3 else "",
                    "part_dispatch_image4": claim_status.part_dispatch_image4.url if claim_status.part_dispatch_image4 else "",
                    "part_dispatch_image5": claim_status.part_dispatch_image5.url if claim_status.part_dispatch_image5 else "",
                    # add other fields as needed
                }
                return JsonResponse({"claim_status_data": data})

    else:
        logout(request)
        messages.error(request, "Unauthorised Access !!")
        return redirect('user_login')    

#############################################################

@login_active_user_required
def check_claim_no_exist(request):
    if request.user.user_profile.user_designation.designation == "claim_manager":
        if request.method == "GET":

            try:
                if request.GET.get('ac_year') and request.GET.get('claim_no'):
                    ac_year = Accounting_Year.objects.filter(ac_year = request.GET.get('ac_year')).first()
                    if ac_year:
                        claim_no_exists = Customer_Information_Report.objects.filter(claim_status__claim_no=request.GET.get('claim_no').lstrip('0'), cir_ac_year = ac_year).first()
                    else:
                        return JsonResponse({"error_msg":"Error in checking claim no !!"})
                else:
                    return JsonResponse({"error_msg":"Error in checking claim no !!"})
                

                if claim_no_exists:
                    if claim_no_exists.job_no != request.GET.get('job_no'):
                        return JsonResponse({"claim_no_exists":"yes", "existing_job_no":claim_no_exists.job_no})
                    else:
                        return JsonResponse({"claim_no_exists":"no"})    
                else:
                    return JsonResponse({"claim_no_exists":"no"})
                
            except:
                return JsonResponse({"error_msg":"Error in checking claim no !!"})
                
    else:
        logout(request)
        messages.error(request, "Unauthorised Access !!")
        return redirect('user_login')


########################################################################################################################################

@login_active_user_required
def my_profile(request):
    if request.user.user_profile.user_designation.designation == "claim_manager":
        user = request.user
        user_profile = user.user_profile  # Or however you access the profile

        if request.method == 'POST':
            form = MyProfileForm(request.POST)
            if form.is_valid():
                # Save logic below!
                user.first_name = form.cleaned_data['user_name'].lower()
                user.save()
                user_profile.mobile_no = form.cleaned_data['mobile_no']
                user_profile.save()
                messages.success(request, "Profile saved successfully")
            
            else:        
                messages.error(request, "Error !!")

            return redirect('cm_my_profile')
    
            
        else:
            # Initial data population
            form = MyProfileForm(initial={
                'user_name': user.first_name,
                'mobile_no': user_profile.mobile_no
            })

            return render(request, 'claim_manager/cm_my_profile.html', {'mform':form})
        
    else:
        logout(request)
        messages.error(request, "Unauthorised Access !!")
        return redirect('user_login')    


########################################################################################################################################

@login_active_user_required
def change_password(request):
    if request.user.user_profile.user_designation.designation == "claim_manager":
            
        if request.method == "POST":    
            pwd_form = Manual_Password_Change_Form(user=request.user, data=request.POST)
            if pwd_form.is_valid():
                pwd_form.save()
                logout(request)
                messages.success(request, "Password changed, Login again with the new Password")
                return redirect('user_login')    

            else:
                messages.error(request, "Old Password is incorrect !!")
                return redirect('cm_change_password')    

        else:
            pwd_form = Manual_Password_Change_Form(user=request.user)
            return render(request, 'claim_manager/cm_change_pwd.html', {'pwd_form':pwd_form})
    

    else:
        logout(request)
        messages.error(request, "Unauthorised Access !!")
        return redirect('user_login')    
